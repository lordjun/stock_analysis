from __future__ import annotations

import argparse
import datetime as dt
import json
import math
import re
import sys
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import akshare as ak
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
import py_mini_racer
import requests
from bs4 import BeautifulSoup
from io import StringIO
from akshare.datasets import get_ths_js
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    import mplfinance as mpf
except Exception:  # pragma: no cover - optional fallback
    mpf = None


REPORT_TITLE = "A股板块龙头日报"
DISCLAIMER = "自动生成，仅供盘后复盘参考，不构成投资建议。"
EASTMONEY_CLIST_URL = "https://push2.eastmoney.com/api/qt/clist/get"
EASTMONEY_UT = "bd1d9ddb04089700cf9c27f6f7426281"
EASTMONEY_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
    "Referer": "https://quote.eastmoney.com/",
}


def ths_headers() -> dict[str, str]:
    with open(get_ths_js("ths.js"), encoding="utf-8") as file:
        js_content = file.read()
    js_code = py_mini_racer.MiniRacer()
    js_code.eval(js_content)
    v_code = js_code.call("v")
    return {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Referer": "http://q.10jqka.com.cn/",
        "Cookie": f"v={v_code}",
    }


@dataclass
class Board:
    name: str
    source: str
    pct_chg: float
    raw: dict
    metric_value: float | None = None
    metric_label: str | None = None


@dataclass
class Leader:
    code: str
    name: str
    pct_chg: float
    board: Board
    raw: dict


def normalize_pct(value) -> float:
    if value is None:
        return float("nan")
    if isinstance(value, str):
        value = value.replace("%", "").replace(",", "").strip()
    try:
        return float(value)
    except Exception:
        return float("nan")


def first_existing(row: pd.Series | dict, names: Iterable[str], default=None):
    for name in names:
        if name in row and pd.notna(row[name]):
            return row[name]
    return default


def clean_stock_code(code: str) -> str:
    text = str(code).strip()
    match = re.search(r"(\d{6})", text)
    return match.group(1) if match else text


def safe_filename_part(text: str) -> str:
    return re.sub(r'[<>:"/\\|?*\s]+', "_", str(text)).strip("_") or "unknown"


def stock_market_prefix(code: str) -> str:
    if code.startswith(("4", "8", "920")):
        return "bj"
    if code.startswith(("6", "9")):
        return "sh"
    return "sz"


def eastmoney_clist(fs: str, fields: str, fid: str = "f3", page_size: int = 5) -> pd.DataFrame:
    """Fetch Eastmoney clist data with small pages to avoid flaky large responses."""
    rows: list[dict] = []
    page = 1
    total = None

    while total is None or len(rows) < total:
        params = {
            "pn": str(page),
            "pz": str(page_size),
            "po": "1",
            "np": "1",
            "ut": EASTMONEY_UT,
            "fltt": "2",
            "invt": "2",
            "fid": fid,
            "fs": fs,
            "fields": fields,
        }

        last_error = None
        for attempt in range(4):
            try:
                session = requests.Session()
                session.trust_env = False
                response = session.get(
                    EASTMONEY_CLIST_URL,
                    params=params,
                    headers=EASTMONEY_HEADERS,
                    timeout=20,
                )
                response.raise_for_status()
                data = response.json().get("data") or {}
                diff = data.get("diff") or []
                total = int(data.get("total") or len(rows) + len(diff))
                rows.extend(diff)
                break
            except Exception as exc:
                last_error = exc
                time.sleep(1.0 + attempt)
        else:
            raise RuntimeError(f"东方财富行情接口请求失败: {last_error}")

        if not diff:
            break
        page += 1
        time.sleep(0.2)

    df = pd.DataFrame(rows)
    if "f3" in df:
        df["f3"] = pd.to_numeric(df["f3"], errors="coerce")
        df.sort_values("f3", ascending=False, inplace=True, ignore_index=True)
    return df


def fetch_boards_em(source: str) -> list[Board]:
    frames: list[tuple[str, pd.DataFrame]] = []
    board_fields = (
        "f2,f3,f4,f8,f12,f14,f20,f104,f105,f128,f136"
    )
    if source in {"industry", "both"}:
        df = eastmoney_clist("m:90 t:2 f:!50", board_fields, fid="f3")
        frames.append(
            (
                "industry",
                df.rename(
                    columns={
                        "f14": "板块名称",
                        "f12": "板块代码",
                        "f3": "涨跌幅",
                        "f2": "最新价",
                        "f4": "涨跌额",
                        "f8": "换手率",
                        "f20": "总市值",
                        "f104": "上涨家数",
                        "f105": "下跌家数",
                        "f128": "领涨股票",
                        "f136": "领涨股票-涨跌幅",
                    }
                ),
            )
        )
    if source in {"concept", "both"}:
        df = eastmoney_clist("m:90 t:3 f:!50", board_fields, fid="f3")
        frames.append(
            (
                "concept",
                df.rename(
                    columns={
                        "f14": "板块名称",
                        "f12": "板块代码",
                        "f3": "涨跌幅",
                        "f2": "最新价",
                        "f4": "涨跌额",
                        "f8": "换手率",
                        "f20": "总市值",
                        "f104": "上涨家数",
                        "f105": "下跌家数",
                        "f128": "领涨股票",
                        "f136": "领涨股票-涨跌幅",
                    }
                ),
            )
        )

    boards: list[Board] = []
    for source_name, df in frames:
        for _, row in df.iterrows():
            name = first_existing(row, ["板块名称", "名称", "行业名称"])
            pct = first_existing(row, ["涨跌幅", "涨跌幅%", "今日涨跌幅", "涨幅"])
            if name:
                boards.append(Board(str(name), source_name, normalize_pct(pct), row.to_dict()))

    return sorted(
        [board for board in boards if not math.isnan(board.pct_chg)],
        key=lambda item: item.pct_chg,
        reverse=True,
    )


def fetch_ths_board_summary(path: str) -> pd.DataFrame:
    """Fetch THS board ranking pages sorted by daily percentage change."""
    headers = ths_headers()
    url = f"http://q.10jqka.com.cn/{path}/index/field/199112/order/desc/page/1/ajax/1/"
    response = requests.get(url, headers=headers, timeout=20)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, features="lxml")
    page_info = soup.find(name="span", attrs={"class": "page_info"})
    page_num = int(page_info.text.split("/")[1]) if page_info else 1

    frames: list[pd.DataFrame] = []
    for page in range(1, page_num + 1):
        page_url = f"http://q.10jqka.com.cn/{path}/index/field/199112/order/desc/page/{page}/ajax/1/"
        response = requests.get(page_url, headers=headers, timeout=20)
        if response.status_code != 200:
            if frames:
                break
            response.raise_for_status()
        try:
            frames.append(pd.read_html(StringIO(response.text))[0])
        except ValueError:
            break
        time.sleep(0.1)

    return pd.concat(frames, ignore_index=True) if frames else pd.DataFrame()


def normalize_ths_board_summary(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df

    if len(df.columns) >= 12:
        df = df.iloc[:, :12].copy()
        df.columns = [
            "序号",
            "板块",
            "涨跌幅",
            "总成交量",
            "总成交额",
            "净流入",
            "上涨家数",
            "下跌家数",
            "均价",
            "领涨股",
            "领涨股-最新价",
            "领涨股-涨跌幅",
        ]
    elif "板块" not in df.columns:
        df = df.rename(columns={df.columns[1]: "板块", df.columns[2]: "涨跌幅"})

    for column in ["涨跌幅", "总成交量", "总成交额", "净流入", "上涨家数", "下跌家数", "均价", "领涨股-最新价", "领涨股-涨跌幅"]:
        if column in df:
            df[column] = pd.to_numeric(df[column], errors="coerce")
    return df


def fetch_ths_concept_boards_mobile(limit: int = 15) -> list[Board]:
    url = f"https://d.10jqka.com.cn/v2/blocksrank/8855/199112/d{limit}.js"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Referer": "https://m.10jqka.com.cn/hq/rank/#concept",
    }
    response = requests.get(url, headers=headers, timeout=20)
    response.raise_for_status()
    response.encoding = "utf-8"
    text = response.text
    payload = json.loads(text[text.find("(") + 1 : text.rfind(")")])

    name_map = ak.stock_board_concept_name_ths()
    code_by_name = {str(row["name"]).strip(): row["code"] for _, row in name_map.iterrows()}

    boards: list[Board] = []
    for item in payload.get("items", []):
        name = str(item.get("55") or "").strip()
        pct = normalize_pct(item.get("199112"))
        if not name or math.isnan(pct):
            continue
        raw = {
            "板块": name,
            "板块代码": code_by_name.get(name),
            "移动端板块代码": item.get("5"),
            "涨跌幅": pct,
            "领涨股": item.get("275"),
            "_provider": "ths",
        }
        boards.append(Board(str(name), "concept", pct, raw))
    return boards


def fetch_fine_concept_boards(limit: int = 50) -> list[Board]:
    df = ak.stock_fund_flow_concept(symbol="即时")
    name_map = ak.stock_board_concept_name_ths()
    code_by_name = {str(row["name"]).strip(): row["code"] for _, row in name_map.iterrows()}

    boards: list[Board] = []
    for _, row in df.head(limit).iterrows():
        name = str(first_existing(row, ["行业", "板块名称", "名称"], "")).strip()
        pct = normalize_pct(first_existing(row, ["行业-涨跌幅", "涨跌幅"], float("nan")))
        if not name or math.isnan(pct):
            continue
        raw = row.to_dict()
        raw["板块"] = name
        raw["板块代码"] = code_by_name.get(name)
        raw["领涨股"] = first_existing(row, ["领涨股"])
        raw["_provider"] = "ths"
        raw["_ranking_provider"] = "eastmoney_fund_flow_concept"
        boards.append(Board(name, "concept", pct, raw))
    return sorted(boards, key=lambda item: item.pct_chg, reverse=True)


def fetch_top_fund_flow_boards(limit: int = 3) -> list[Board]:
    df = ak.stock_fund_flow_concept(symbol="即时").copy()
    df["净额"] = pd.to_numeric(df["净额"], errors="coerce")
    df.sort_values("净额", ascending=False, inplace=True, ignore_index=True)

    name_map = ak.stock_board_concept_name_ths()
    code_by_name = {str(row["name"]).strip(): row["code"] for _, row in name_map.iterrows()}
    try:
        spot = ak.stock_sector_spot(indicator="概念")
        label_by_name = {str(row["板块"]).strip(): row["label"] for _, row in spot.iterrows()}
    except Exception:
        label_by_name = {}

    boards: list[Board] = []
    for _, row in df.head(limit).iterrows():
        name = str(row["行业"]).strip()
        pct = normalize_pct(row.get("行业-涨跌幅"))
        net = normalize_pct(row.get("净额"))
        raw = row.to_dict()
        raw["板块"] = name
        raw["板块代码"] = code_by_name.get(name)
        raw["_sina_label"] = label_by_name.get(name)
        raw["_provider"] = "ths"
        boards.append(Board(name, "concept", pct, raw, metric_value=net, metric_label="主力净流入"))
    return boards


def fetch_top_volume_boards(limit: int = 3) -> list[Board]:
    df = ak.stock_sector_spot(indicator="概念").copy()
    df["总成交量"] = pd.to_numeric(df["总成交量"], errors="coerce")
    df["涨跌幅"] = pd.to_numeric(df["涨跌幅"], errors="coerce")
    df.sort_values("总成交量", ascending=False, inplace=True, ignore_index=True)

    boards: list[Board] = []
    for _, row in df.head(limit).iterrows():
        name = str(row["板块"]).strip()
        volume = normalize_pct(row.get("总成交量"))
        pct = normalize_pct(row.get("涨跌幅"))
        raw = row.to_dict()
        raw["_sina_label"] = row.get("label")
        raw["_provider"] = "sina"
        boards.append(Board(name, "concept", pct, raw, metric_value=volume, metric_label="成交量"))
    return boards


def fetch_ths_realtime_stocks(codes: Iterable[str]) -> dict[str, dict]:
    clean_codes = [clean_stock_code(code) for code in codes if code]
    if not clean_codes:
        return {}
    url = "https://d.10jqka.com.cn/real/hsa/" + ",".join(clean_codes)
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Referer": "https://m.10jqka.com.cn/hq/rank/",
    }
    response = requests.get(url, headers=headers, timeout=20)
    response.raise_for_status()
    response.encoding = "utf-8"
    text = response.text
    return json.loads(text[text.find("(") + 1 : text.rfind(")")])


def fetch_boards_ths(source: str) -> list[Board]:
    frames: list[tuple[str, pd.DataFrame, pd.DataFrame]] = []
    boards: list[Board] = []

    if source in {"industry", "both"}:
        frames.append(
            (
                "industry",
                ak.stock_board_industry_name_ths(),
                normalize_ths_board_summary(fetch_ths_board_summary("thshy")),
            )
        )

    if source == "concept":
        try:
            boards.extend(fetch_fine_concept_boards())
        except Exception:
            boards.extend(fetch_ths_concept_boards_mobile())
    elif source == "both":
        boards.extend(fetch_ths_concept_boards_mobile())

    for source_name, name_map, df in frames:
        code_by_name = dict(zip(name_map["name"], name_map["code"]))
        for _, row in df.iterrows():
            name = first_existing(row, ["板块", "板块名称", "名称"])
            pct = first_existing(row, ["涨跌幅", "涨跌幅%", "今日涨跌幅", "涨幅"])
            if not name:
                continue
            raw = row.to_dict()
            raw["板块代码"] = code_by_name.get(str(name))
            raw["_provider"] = "ths"
            boards.append(Board(str(name), source_name, normalize_pct(pct), raw))

    source_priority = {"industry": 0, "concept": 1}
    return sorted(
        [board for board in boards if not math.isnan(board.pct_chg)],
        key=lambda item: (-item.pct_chg, source_priority.get(item.source, 9), item.name),
    )


def fetch_boards(source: str, data_source: str) -> list[Board]:
    if data_source == "ths":
        return fetch_boards_ths(source)
    return fetch_boards_em(source)


def fetch_board_cons(board: Board) -> pd.DataFrame:
    if board.raw.get("_sina_label"):
        df = ak.stock_sector_detail(str(board.raw["_sina_label"]))
        return df.rename(
            columns={
                "code": "代码",
                "name": "名称",
                "trade": "最新价",
                "changepercent": "涨跌幅",
                "pricechange": "涨跌额",
                "turnoverratio": "换手率",
            }
        )

    if board.raw.get("_provider") == "ths":
        board_code = first_existing(board.raw, ["板块代码"])
        if not board_code:
            raise RuntimeError(f"没有找到同花顺板块代码: {board.name}")

        path = "thshy" if board.source == "industry" else "gn"
        url = f"http://q.10jqka.com.cn/{path}/detail/code/{board_code}/"
        response = requests.get(url, headers=ths_headers(), timeout=20)
        response.raise_for_status()
        response.encoding = response.apparent_encoding
        try:
            df = pd.read_html(StringIO(response.text))[0]
        except ValueError:
            raise RuntimeError(f"同花顺板块成分股表解析失败: {board.name}")

        if "代码" in df and not df.empty:
            df = df[df["代码"].astype(str).str.contains(r"\d{6}", regex=True, na=False)]
        return df.rename(
            columns={
                "涨跌幅(%)": "涨跌幅",
                "现价": "最新价",
                "涨跌": "涨跌额",
                "换手(%)": "换手率",
                "振幅(%)": "振幅",
            }
        )

    board_code = first_existing(board.raw, ["板块代码", "f12"])
    if not board_code:
        if board.source == "industry":
            return ak.stock_board_industry_cons_em(symbol=board.name)
        return ak.stock_board_concept_cons_em(symbol=board.name)

    cons_fields = (
        "f2,f3,f4,f5,f6,f7,f8,f9,f10,f12,f14,f15,f16,f17,f18,f20,f21,f23,"
        "f24,f25,f22,f11,f62,f128,f136,f115,f152,f45"
    )
    df = eastmoney_clist(f"b:{board_code} f:!50", cons_fields, fid="f3")
    return df.rename(
        columns={
            "f12": "代码",
            "f14": "名称",
            "f2": "最新价",
            "f3": "涨跌幅",
            "f4": "涨跌额",
            "f5": "成交量",
            "f6": "成交额",
            "f7": "振幅",
            "f8": "换手率",
            "f9": "市盈率-动态",
            "f15": "最高",
            "f16": "最低",
            "f17": "今开",
            "f18": "昨收",
            "f23": "市净率",
        }
    )


def top_leaders(board: Board, limit: int = 3) -> list[Leader]:
    try:
        df = fetch_board_cons(board)
    except Exception:
        df = pd.DataFrame()
    rows: list[Leader] = []
    for _, row in df.iterrows():
        code = clean_stock_code(first_existing(row, ["代码", "股票代码", "证券代码"], ""))
        name = first_existing(row, ["名称", "股票名称", "证券简称"], code)
        pct = first_existing(row, ["涨跌幅", "涨跌幅%", "涨幅"], float("nan"))
        if code and name:
            rows.append(Leader(code, str(name), normalize_pct(pct), board, row.to_dict()))
    leaders = sorted(
        [item for item in rows if not math.isnan(item.pct_chg)],
        key=lambda item: item.pct_chg,
        reverse=True,
    )[:limit]
    if leaders:
        return leaders

    if board.source == "concept":
        try:
            spot = ak.stock_sector_spot(indicator="概念")
            matched = spot[spot["板块"].astype(str).str.strip() == board.name]
            if not matched.empty:
                label = matched.iloc[0]["label"]
                df = ak.stock_sector_detail(label)
                sina_rows: list[Leader] = []
                for _, row in df.iterrows():
                    code = clean_stock_code(first_existing(row, ["code", "symbol"], ""))
                    name = first_existing(row, ["name"], code)
                    pct = normalize_pct(first_existing(row, ["changepercent"], float("nan")))
                    if code and name:
                        sina_rows.append(Leader(code, str(name), pct, board, row.to_dict()))
                leaders = sorted(
                    [item for item in sina_rows if not math.isnan(item.pct_chg)],
                    key=lambda item: item.pct_chg,
                    reverse=True,
                )[:limit]
                if leaders:
                    return leaders
        except Exception:
            pass

    lead_code = clean_stock_code(board.raw.get("领涨股", ""))
    if board.raw.get("_provider") == "ths" and re.fullmatch(r"\d{6}", lead_code or ""):
        realtime = fetch_ths_realtime_stocks([lead_code]).get(lead_code, {})
        name = realtime.get("name", lead_code)
        pct = normalize_pct(realtime.get("199112"))
        if name and not math.isnan(pct):
            return [Leader(lead_code, str(name), pct, board, realtime)]
    lead_name = first_existing(board.raw, ["领涨股"])
    lead_pct = normalize_pct(first_existing(board.raw, ["领涨股-涨跌幅"], float("nan")))
    if lead_name and not math.isnan(lead_pct):
        return [Leader("-", str(lead_name), lead_pct, board, board.raw)]
    return []


def fetch_stock_hist(code: str, end_date: str, lookback_days: int) -> pd.DataFrame:
    end = dt.datetime.strptime(end_date, "%Y%m%d").date()
    start = end - dt.timedelta(days=max(lookback_days * 2, 120))
    start_text = start.strftime("%Y%m%d")
    end_text = end.strftime("%Y%m%d")
    market_code = stock_market_prefix(code) + code

    try:
        df = ak.stock_zh_a_hist(
            symbol=code,
            period="daily",
            start_date=start_text,
            end_date=end_text,
            adjust="qfq",
        )
        if not df.empty:
            return df.tail(lookback_days).copy()
    except Exception:
        pass

    for fetcher in (
        lambda: ak.stock_zh_a_hist_tx(symbol=market_code, start_date=start_text, end_date=end_text, adjust="qfq"),
        lambda: ak.stock_zh_a_daily(symbol=market_code, start_date=start_text, end_date=end_text, adjust="qfq"),
    ):
        try:
            df = fetcher()
            if df.empty:
                continue
            df = df.rename(
                columns={
                    "date": "日期",
                    "open": "开盘",
                    "close": "收盘",
                    "high": "最高",
                    "low": "最低",
                    "volume": "成交量",
                    "amount": "成交量",
                }
            )
            df = df.loc[:, ~df.columns.duplicated()]
            return df.tail(lookback_days).copy()
        except Exception:
            continue

    try:
        df = fetch_stock_hist_tx_direct(market_code, start_text, end_text, adjust="qfq")
        if not df.empty:
            df = df.rename(
                columns={
                    "date": "日期",
                    "open": "开盘",
                    "close": "收盘",
                    "high": "最高",
                    "low": "最低",
                    "amount": "成交量",
                }
            )
            return df.tail(lookback_days).copy()
    except Exception:
        pass

    return pd.DataFrame()


def fetch_stock_hist_tx_direct(symbol: str, start_date: str, end_date: str, adjust: str = "qfq") -> pd.DataFrame:
    url = "https://proxy.finance.qq.com/ifzqgtimg/appstock/app/newfqkline/get"
    frames: list[pd.DataFrame] = []
    start_year = int(start_date[:4])
    end_year = int(end_date[:4])
    key_order = ("day", "qfqday", "hfqday")

    for year in range(start_year, end_year + 1):
        params = {
            "_var": f"kline_day{adjust}{year}",
            "param": f"{symbol},day,{year}-01-01,{year + 1}-12-31,640,{adjust}",
            "r": "0.8205512681390605",
        }
        response = requests.get(url, params=params, timeout=20)
        response.raise_for_status()
        data_text = response.text
        payload = json.loads(data_text[data_text.find("={") + 1 :])
        stock_data = payload.get("data", {}).get(symbol, {})
        rows = []
        for key in key_order:
            if key in stock_data:
                rows = stock_data[key]
                break
        if rows:
            frames.append(pd.DataFrame(rows))

    if not frames:
        return pd.DataFrame()

    df = pd.concat(frames, ignore_index=True).iloc[:, :6]
    df.columns = ["date", "open", "close", "high", "low", "amount"]
    df["date"] = pd.to_datetime(df["date"], errors="coerce").dt.date
    for column in ["open", "close", "high", "low", "amount"]:
        df[column] = pd.to_numeric(df[column], errors="coerce")
    df = df.dropna(subset=["date"]).drop_duplicates(ignore_index=True)
    df.index = pd.to_datetime(df["date"], errors="coerce")
    df = df.sort_index()[start_date:end_date].reset_index(drop=True)
    return df


def fetch_stock_news(code: str, max_rows: int = 30) -> pd.DataFrame:
    try:
        return ak.stock_news_em(symbol=code).head(max_rows)
    except Exception:
        return pd.DataFrame()


def recent_news_titles(code: str, news_days: int) -> list[str]:
    df = fetch_stock_news(code)
    if df.empty:
        return []

    cutoff = dt.datetime.now() - dt.timedelta(days=news_days)
    titles: list[str] = []
    for _, row in df.iterrows():
        title = first_existing(row, ["新闻标题", "标题", "title"])
        if not title:
            continue
        date_value = first_existing(row, ["发布时间", "时间", "日期"])
        if date_value is not None:
            parsed = pd.to_datetime(date_value, errors="coerce")
            if pd.notna(parsed) and parsed.to_pydatetime().replace(tzinfo=None) < cutoff:
                continue
        titles.append(str(title))
    return titles[:8]


KEYWORDS = {
    "政策利好": ["政策", "方案", "规划", "意见", "通知", "改革", "国资", "补贴", "税", "会议"],
    "产业趋势": ["AI", "人工智能", "算力", "机器人", "新能源", "半导体", "芯片", "低空", "数据", "国产替代"],
    "业绩改善": ["业绩", "增长", "利润", "订单", "中标", "合同", "预增", "扭亏", "景气"],
    "资金活跃": ["涨停", "主力", "资金", "成交", "放量", "龙虎榜", "北向", "融资"],
    "事件催化": ["发布", "公告", "合作", "并购", "重组", "收购", "新品", "上市", "突破"],
}


def load_reason_notes(path: str | None) -> dict[str, dict]:
    if not path:
        return {}
    note_path = Path(path)
    if not note_path.exists():
        raise FileNotFoundError(f"原因分析文件不存在: {note_path}")
    return json.loads(note_path.read_text(encoding="utf-8"))


def analyze_reasons(
    board: Board,
    leaders: list[Leader],
    news_map: dict[str, list[str]],
    manual_note: dict | None = None,
) -> list[str]:
    if manual_note and manual_note.get("reasons"):
        return [str(item) for item in manual_note["reasons"]][:5]

    corpus = " ".join([board.name] + [leader.name for leader in leaders])
    for titles in news_map.values():
        corpus += " " + " ".join(titles)

    reasons: list[str] = []
    for label, words in KEYWORDS.items():
        hits = [word for word in words if word.lower() in corpus.lower()]
        if hits:
            reasons.append(f"{label}: 相关新闻/名称中出现 {', '.join(hits[:5])} 等线索。")

    if not reasons:
        reasons.append("市场表现: 板块及龙头股涨幅靠前，短线资金关注度提升。")

    leading_names = "、".join(f"{leader.name}({leader.pct_chg:.2f}%)" for leader in leaders)
    reasons.insert(0, f"领涨结构: {leading_names} 对板块涨幅贡献明显。")
    return reasons[:4]


def fit_lines(lines: Iterable[str], max_chars: int = 180, max_lines: int = 3) -> list[str]:
    result: list[str] = []
    used = 0
    for line in lines:
        text = re.sub(r"\s+", " ", str(line)).strip()
        if not text:
            continue
        remaining = max_chars - used
        if remaining <= 0:
            break
        if len(text) > remaining:
            text = text[: max(0, remaining - 1)].rstrip("，,。；; ") + "..."
        result.append(text)
        used += len(text)
        if len(result) >= max_lines:
            break
    return result


def ensure_chinese_font() -> None:
    plt.rcParams["font.sans-serif"] = [
        "Microsoft YaHei",
        "SimHei",
        "Arial Unicode MS",
        "DejaVu Sans",
    ]
    plt.rcParams["axes.unicode_minus"] = False


def save_kline_chart(leader: Leader, end_date: str, output_dir: Path, lookback_days: int) -> Path | None:
    try:
        hist = fetch_stock_hist(leader.code, end_date, lookback_days)
    except Exception:
        return None
    if hist.empty:
        return None

    ensure_chinese_font()
    chart_path = output_dir / f"{leader.code}_{safe_filename_part(leader.name)}_kline.png"
    chart_path.parent.mkdir(parents=True, exist_ok=True)

    required = {"日期", "开盘", "收盘", "最高", "最低", "成交量"}
    if mpf is not None and required.issubset(set(hist.columns)):
        data = hist.rename(
            columns={
                "日期": "Date",
                "开盘": "Open",
                "收盘": "Close",
                "最高": "High",
                "最低": "Low",
                "成交量": "Volume",
            }
        )
        data["Date"] = pd.to_datetime(data["Date"])
        data = data.set_index("Date")[["Open", "High", "Low", "Close", "Volume"]]
        market_colors = mpf.make_marketcolors(up="#d62728", down="#2ca02c", edge="inherit", wick="inherit", volume="inherit")
        style = mpf.make_mpf_style(base_mpf_style="yahoo", marketcolors=market_colors, rc={"font.sans-serif": "Microsoft YaHei"})
        mpf.plot(
            data,
            type="candle",
            mav=(5, 10, 20),
            volume=True,
            style=style,
            title=f"{leader.name}({leader.code}) 近{lookback_days}日K线",
            savefig=dict(fname=str(chart_path), dpi=150, bbox_inches="tight"),
        )
    else:
        fig, ax = plt.subplots(figsize=(8, 3.8), dpi=150)
        x = range(len(hist))
        ax.plot(list(x), hist["收盘"], color="#1f77b4", linewidth=1.6, label="收盘")
        if "开盘" in hist:
            ax.plot(list(x), hist["开盘"], color="#8c564b", linewidth=1.0, alpha=0.7, label="开盘")
        ax.set_title(f"{leader.name}({leader.code}) 近{lookback_days}日走势")
        ax.grid(alpha=0.25)
        ax.legend()
        fig.tight_layout()
        fig.savefig(chart_path)
        plt.close(fig)
    return chart_path


BG_COLOR = RGBColor(246, 248, 251)
NAV_COLOR = RGBColor(18, 42, 68)
TEXT_COLOR = RGBColor(31, 41, 55)
MUTED_COLOR = RGBColor(100, 116, 139)
ACCENT_RED = RGBColor(190, 49, 49)
ACCENT_GOLD = RGBColor(197, 139, 45)
CARD_BORDER = RGBColor(218, 226, 236)
CARD_FILL = RGBColor(255, 255, 255)


def add_slide_background(slide, section: str | None = None) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    band.fill.solid()
    band.fill.fore_color.rgb = NAV_COLOR
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.95), Inches(0.32), Inches(0.18), Inches(6.85))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_RED
    accent.line.fill.background()

    if section:
        add_textbox(
            slide,
            Inches(0.58),
            Inches(0.06),
            Inches(5.0),
            Inches(0.18),
            [section],
            font_size=7,
            color=RGBColor(232, 238, 245),
            bold=True,
        )


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    title = slide.shapes.add_textbox(Inches(0.58), Inches(0.55), Inches(8.4), Inches(0.55))
    frame = title.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.06), Inches(10.5), Inches(0.28))
        sub_frame = sub.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED_COLOR


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: Iterable[str],
    font_size: int = 14,
    color: RGBColor = TEXT_COLOR,
    bold: bool = False,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = str(line)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(5)


def add_hyperlink_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: Iterable[str],
    font_size: int = 10,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.clear()
    url_pattern = re.compile(r"https?://\S+")

    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.space_after = Pt(8)
        text = str(line)
        pos = 0
        for match in url_pattern.finditer(text):
            if match.start() > pos:
                run = p.add_run()
                run.text = text[pos : match.start()]
                run.font.size = Pt(font_size)
                run.font.color.rgb = TEXT_COLOR
            url = match.group(0).rstrip("，,。.;；")
            run = p.add_run()
            run.text = url
            run.hyperlink.address = url
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(37, 99, 235)
            run.font.underline = True
            pos = match.start() + len(match.group(0))
        if pos < len(text):
            run = p.add_run()
            run.text = text[pos:]
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_COLOR


def add_card(slide, left, top, width, height, fill: RGBColor = CARD_FILL, line: RGBColor = CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_metric_card(slide, left, top, width, height, rank: int, name: str, pct_chg: float, source: str) -> None:
    add_card(slide, left, top, width, height)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18), top + Inches(0.18), Inches(0.46), Inches(0.46))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_RED if rank == 1 else ACCENT_GOLD
    badge.line.fill.background()
    frame = badge.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = str(rank)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    add_textbox(slide, left + Inches(0.76), top + Inches(0.17), width - Inches(0.95), Inches(0.38), [name], font_size=15, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.76), Inches(1.35), Inches(0.42), [f"{pct_chg:.2f}%"], font_size=21, color=ACCENT_RED, bold=True)
    add_textbox(slide, left + Inches(1.52), top + Inches(0.88), width - Inches(1.75), Inches(0.25), [source], font_size=9, color=MUTED_COLOR)


def add_footer(slide, report_date: str) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(7.08), Inches(12.0), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()
    add_textbox(slide, Inches(0.58), Inches(7.11), Inches(6.2), Inches(0.2), [DISCLAIMER], font_size=7, color=MUTED_COLOR)
    add_textbox(slide, Inches(10.5), Inches(7.11), Inches(2.1), Inches(0.2), [report_date], font_size=7, color=MUTED_COLOR)


def format_metric(value: float | None, metric_label: str | None) -> str:
    if value is None or math.isnan(value):
        return "-"
    if metric_label == "主力净流入":
        return f"{value:.2f}亿"
    if metric_label == "成交量":
        if abs(value) >= 100000000:
            return f"{value / 100000000:.2f}亿手"
        if abs(value) >= 10000:
            return f"{value / 10000:.2f}万手"
        return f"{value:.0f}手"
    return f"{value:.2f}"


def add_table(slide, left, top, width, height, rows: list[list[str]], header_fill=NAV_COLOR) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    for row_idx, row in enumerate(rows):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(text)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.color.rgb = TEXT_COLOR
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(243, 246, 250)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)


def build_ppt(
    report_date: str,
    boards: list[Board],
    ranking_sections: dict[str, list[Board]],
    leaders_by_board: dict[str, list[Leader]],
    reasons_by_board: dict[str, list[str]],
    notes_by_board: dict[str, dict],
    charts: dict[str, Path | None],
    output_path: Path,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank_layout)
    add_slide_background(title_slide, "A-SHARE DAILY SECTOR REVIEW")
    cover_band = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.32), Inches(13.333), Inches(2.6))
    cover_band.fill.solid()
    cover_band.fill.fore_color.rgb = NAV_COLOR
    cover_band.line.fill.background()
    add_textbox(title_slide, Inches(0.72), Inches(0.86), Inches(8.6), Inches(0.72), [REPORT_TITLE], font_size=31, color=RGBColor(255, 255, 255), bold=True)
    add_textbox(title_slide, Inches(0.78), Inches(1.72), Inches(8.6), Inches(0.34), [f"{report_date} 收盘后复盘 | 严格按全市场板块涨幅排序"], font_size=13, color=RGBColor(218, 226, 236))
    add_textbox(title_slide, Inches(0.78), Inches(2.2), Inches(8.6), Inches(0.28), [DISCLAIMER], font_size=9, color=RGBColor(190, 203, 220))

    for idx, board in enumerate(boards, start=1):
        add_metric_card(
            title_slide,
            Inches(0.72 + (idx - 1) * 4.14),
            Inches(3.38),
            Inches(3.72),
            Inches(1.45),
            idx,
            board.name,
            board.pct_chg,
            "行业板块" if board.source == "industry" else "概念板块",
        )
    add_footer(title_slide, report_date)

    summary = prs.slides.add_slide(blank_layout)
    add_slide_background(summary, "SUMMARY")
    add_title(summary, "涨幅前三板块", "同口径混排，严格保留涨幅前三；细分板块只有在涨幅超过大板块时才进入结果。")
    rows = [["排名", "板块", "口径", "涨跌幅"]]
    for idx, board in enumerate(boards, start=1):
        rows.append([str(idx), board.name, "行业" if board.source == "industry" else "概念", f"{board.pct_chg:.2f}%"])
    add_card(summary, Inches(0.62), Inches(1.52), Inches(5.7), Inches(2.0))
    add_table(summary, Inches(0.82), Inches(1.75), Inches(5.3), Inches(1.45), rows)

    add_card(summary, Inches(6.62), Inches(1.52), Inches(5.72), Inches(2.0))
    add_textbox(summary, Inches(6.88), Inches(1.75), Inches(5.15), Inches(0.32), ["口径说明"], font_size=13, bold=True)
    add_textbox(
        summary,
        Inches(6.88),
        Inches(2.14),
        Inches(5.15),
        Inches(1.0),
        [
            "行业板块与概念板块使用同一数据源口径排序。",
            "如果某个大板块涨幅更高，则保留大板块；不为追求细分替换排名。",
            "每个入选板块继续列出板块内涨幅前三个股。",
        ],
        font_size=10,
        color=MUTED_COLOR,
    )

    y = 3.9
    for idx, board in enumerate(boards, start=1):
        add_card(summary, Inches(0.62), Inches(y), Inches(11.72), Inches(0.82))
        add_textbox(summary, Inches(0.86), Inches(y + 0.13), Inches(1.0), Inches(0.28), [f"#{idx}"], font_size=15, color=ACCENT_RED, bold=True)
        add_textbox(summary, Inches(1.55), Inches(y + 0.13), Inches(1.8), Inches(0.28), [board.name], font_size=14, bold=True)
        add_textbox(summary, Inches(3.2), Inches(y + 0.15), Inches(1.05), Inches(0.24), [f"{board.pct_chg:.2f}%"], font_size=13, color=ACCENT_RED, bold=True)
        add_textbox(summary, Inches(4.25), Inches(y + 0.12), Inches(7.8), Inches(0.42), ["；".join(reasons_by_board[board.name])], font_size=9, color=MUTED_COLOR)
        y += 0.92
    add_footer(summary, report_date)

    for section_title, section_boards in ranking_sections.items():
        slide = prs.slides.add_slide(blank_layout)
        add_slide_background(slide, "ADDITIONAL RANKING")
        add_title(slide, section_title, "按对应指标排序，列出前三板块及各板块内涨幅前三个股。")
        rows = [["排名", "板块", "指标", "涨幅前三个股"]]
        for idx, board in enumerate(section_boards, start=1):
            leaders = leaders_by_board.get(board.name, [])
            leader_text = "；".join(f"{item.name} {item.pct_chg:.2f}%" for item in leaders[:3]) or "-"
            rows.append(
                [
                    str(idx),
                    board.name,
                    format_metric(board.metric_value, board.metric_label),
                    leader_text,
                ]
            )
        add_card(slide, Inches(0.68), Inches(1.55), Inches(11.75), Inches(4.75))
        add_table(slide, Inches(0.92), Inches(1.9), Inches(11.25), Inches(2.6), rows)
        add_textbox(
            slide,
            Inches(0.95),
            Inches(4.85),
            Inches(10.95),
            Inches(0.65),
            ["说明：该页不改变日报主排名；主排名仍严格按全市场板块涨幅前三输出。"],
            font_size=9,
            color=MUTED_COLOR,
        )
        add_footer(slide, report_date)

    for board in boards:
        leaders = leaders_by_board[board.name]
        slide = prs.slides.add_slide(blank_layout)
        add_slide_background(slide, "BOARD DETAIL")
        add_title(slide, f"{board.name} 龙头股", f"{'行业' if board.source == 'industry' else '概念'}板块 | 板块涨跌幅 {board.pct_chg:.2f}%")

        leader_rows = [["排名", "代码", "名称", "涨跌幅"]]
        for idx, leader in enumerate(leaders, start=1):
            leader_rows.append([str(idx), leader.code, leader.name, f"{leader.pct_chg:.2f}%"])
        add_card(slide, Inches(0.58), Inches(1.55), Inches(5.1), Inches(1.65))
        add_table(slide, Inches(0.78), Inches(1.78), Inches(4.7), Inches(1.2), leader_rows)
        add_card(slide, Inches(5.96), Inches(1.55), Inches(6.32), Inches(1.82))
        add_textbox(
            slide,
            Inches(6.18),
            Inches(1.72),
            Inches(5.88),
            Inches(1.44),
            fit_lines(reasons_by_board[board.name], max_chars=165, max_lines=3),
            font_size=8,
            color=MUTED_COLOR,
        )

        chart_positions = [(0.62, 3.96), (4.58, 3.96), (8.54, 3.96)]
        for leader, (left, top) in zip(leaders, chart_positions):
            chart = charts.get(leader.code)
            add_card(slide, Inches(left), Inches(top - 0.42), Inches(3.42), Inches(2.92))
            add_textbox(slide, Inches(left + 0.14), Inches(top - 0.3), Inches(3.05), Inches(0.24), [f"{leader.name}({leader.code})"], font_size=10, bold=True)
            if chart and chart.exists():
                slide.shapes.add_picture(str(chart), Inches(left + 0.13), Inches(top + 0.02), width=Inches(3.16), height=Inches(2.26))
            else:
                add_textbox(slide, Inches(left + 0.2), Inches(top + 0.72), Inches(3.0), Inches(0.6), ["K线图生成失败"], font_size=12, color=MUTED_COLOR)

        news_lines = []
        for leader in leaders:
            titles = recent_news_titles(leader.code, 5)[:2]
            if titles:
                news_lines.append(f"{leader.name}: " + "；".join(titles))
        if news_lines:
            add_textbox(
                slide,
                Inches(0.72),
                Inches(6.72),
                Inches(11.65),
                Inches(0.24),
                fit_lines(news_lines, max_chars=130, max_lines=1),
                font_size=7,
                color=MUTED_COLOR,
            )
        add_footer(slide, report_date)

    source_lines: list[str] = []
    for board in boards:
        note = notes_by_board.get(board.name, {})
        for source in note.get("sources", []):
            source_lines.append(f"{board.name}: {source}")
    if source_lines:
        slide = prs.slides.add_slide(blank_layout)
        add_slide_background(slide, "SOURCES")
        add_title(slide, "参考信息")
        add_card(slide, Inches(0.68), Inches(1.5), Inches(11.75), Inches(5.35))
        add_hyperlink_textbox(slide, Inches(0.95), Inches(1.78), Inches(11.2), Inches(4.8), source_lines, font_size=10)
        add_footer(slide, report_date)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="生成A股板块龙头PPT日报")
    parser.add_argument("--date", default=dt.date.today().strftime("%Y%m%d"), help="交易日期，格式YYYYMMDD")
    parser.add_argument("--board-source", choices=["industry", "concept", "both"], default="both")
    parser.add_argument("--data-source", choices=["ths", "em"], default="ths", help="板块数据源：ths=同花顺，em=东方财富")
    parser.add_argument("--output-dir", default="reports")
    parser.add_argument("--lookback-days", type=int, default=90)
    parser.add_argument("--news-days", type=int, default=14)
    parser.add_argument("--reason-file", default=None, help="可选 JSON 文件，用于覆盖/增强板块上涨原因分析")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    output_dir = Path(args.output_dir)
    chart_dir = output_dir / "charts" / args.date
    notes_by_board = load_reason_notes(args.reason_file)

    print("Fetching board ranking...")
    boards = fetch_boards(args.board_source, args.data_source)[:3]
    if not boards:
        raise RuntimeError("没有获取到板块涨幅数据，请确认 Akshare 接口和网络可用。")

    print("Fetching additional rankings...")
    ranking_sections: dict[str, list[Board]] = {
        "主力资金流入最多的3个板块": fetch_top_fund_flow_boards(3),
        "当日成交量最大的3个板块": fetch_top_volume_boards(3),
    }

    leaders_by_board: dict[str, list[Leader]] = {}
    reasons_by_board: dict[str, list[str]] = {}
    charts: dict[str, Path | None] = {}

    all_boards: list[Board] = []
    seen_boards: set[str] = set()
    for board in boards + [item for section in ranking_sections.values() for item in section]:
        if board.name in seen_boards:
            continue
        seen_boards.add(board.name)
        all_boards.append(board)

    for board in all_boards:
        print(f"Fetching leaders for {board.name}...")
        leaders = top_leaders(board, 3)
        leaders_by_board[board.name] = leaders

        if board in boards:
            news_map = {leader.code: recent_news_titles(leader.code, args.news_days) for leader in leaders}
            reasons_by_board[board.name] = analyze_reasons(board, leaders, news_map, notes_by_board.get(board.name))

            for leader in leaders:
                print(f"Drawing K-line chart for {leader.name}({leader.code})...")
                charts[leader.code] = save_kline_chart(leader, args.date, chart_dir, args.lookback_days)

    report_path = output_dir / f"A股板块龙头日报_{args.date}.pptx"
    build_ppt(args.date, boards, ranking_sections, leaders_by_board, reasons_by_board, notes_by_board, charts, report_path)
    print(f"Report saved: {report_path.resolve()}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except KeyboardInterrupt:
        raise SystemExit(130)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise SystemExit(1)
